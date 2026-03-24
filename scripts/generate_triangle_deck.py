from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOC_DIR = ROOT / "doc"
PUBLIC_DIR = ROOT / "frontend" / "public"
OUTPUT = DOC_DIR / "Triangle-Dealership-Ride-Fleet-Loaner-Program-2026-03-24.pptx"
LOGO = PUBLIC_DIR / "ride-logo.png"

PRIMARY = RGBColor(98, 61, 239)
PRIMARY_DARK = RGBColor(32, 27, 76)
PRIMARY_SOFT = RGBColor(241, 236, 255)
TEXT = RGBColor(44, 43, 65)
MUTED = RGBColor(108, 107, 136)
SUCCESS = RGBColor(37, 133, 79)
WARNING = RGBColor(197, 109, 34)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 249, 255)


def add_header_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.25), Inches(12.63), Inches(1.15)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY_SOFT
    band.line.color.rgb = PRIMARY_SOFT

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(7.8), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = PRIMARY_DARK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.67), Inches(0.82), Inches(8.5), Inches(0.28))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.color.rgb = MUTED


def add_footer(slide, text="Ride Fleet | Triangle Dealership discussion | March 24, 2026"):
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.92), Inches(12.0), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_bullets(slide, bullets, left=0.8, top=1.55, width=11.2, height=5.0, font_size=21):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(9)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = bullet
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT


def add_two_column_bullets(
    slide,
    left_title,
    left_bullets,
    right_title,
    right_bullets,
    top=1.6,
):
    def column(x, title, bullets, accent):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(5.85), Inches(4.95)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = PRIMARY_SOFT

        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(top + 0.18), Inches(5.1), Inches(0.3))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = accent

        add_bullets(slide, bullets, left=x + 0.23, top=top + 0.6, width=5.25, height=4.0, font_size=15)

    column(0.55, left_title, left_bullets, PRIMARY_DARK)
    column(6.1, right_title, right_bullets, PRIMARY)


def add_pricing_slide(slide):
    add_header_band(slide, "Recommended pricing strategy", "Competitive against enterprise players while still protecting margin")

    headers = [
        ("Pilot Launch", PRIMARY),
        ("Core Dealership", SUCCESS),
        ("Multi-Rooftop Group", WARNING),
    ]
    bodies = [
        [
            "$0 implementation fee",
            "60-day pilot",
            "1 rooftop",
            "Up to 25 active loaners",
            "$199/month",
        ],
        [
            "1 rooftop",
            "Up to 40 active loaners included",
            "$399/month",
            "$7/additional active unit",
            "Best standard launch package",
        ],
        [
            "2+ rooftops",
            "Central reporting and group controls",
            "$699-$1,499/month starting range",
            "Custom quote by scale and rollout needs",
        ],
    ]
    xs = [0.55, 4.35, 8.15]
    for idx, x in enumerate(xs):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(3.2), Inches(4.85)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = PRIMARY_SOFT

        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(1.9), Inches(1.95), Inches(0.42)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = headers[idx][1]
        pill.line.color.rgb = headers[idx][1]
        tf = pill.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = headers[idx][0]
        run.font.name = "Aptos"
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(255, 255, 255)

        add_bullets(slide, bodies[idx], left=x + 0.2, top=2.5, width=2.75, height=3.6, font_size=15)


def add_table_like_slide(slide):
    add_header_band(slide, "Competitive market view", "Position Ride Fleet between legacy and premium enterprise")
    left = [
        "Bluebird LoanerTrack: legacy dealership footprint, reservation module, retail rental contracts, report-heavy positioning.",
        "Dealerware: modern dealership-first mobile story, strong pitch around faster contracting and utilization.",
        "TSD Mobility: enterprise scale, 16,000+ dealerships cited by Reynolds, broad mobility scope.",
        "ARSLoaner: cloud-native, transparent pricing, aggressive value proposition.",
        "RENTALL: dealership loaner offering with cloud/mobile/payments/reporting language and a custom-sales motion.",
    ]
    right = [
        "Ride Fleet can win with modern mobile UX, accounting packets, issue center, flexible tenant setup, and faster customization.",
        "Strong today: intake, borrower packet, advisor ops, billing control, extensions, swaps, service completion, monthly packets.",
        "Be careful not to overclaim yet: DMS depth, OEM-specific libraries, telematics automation, insurer network integrations.",
        "Best positioning: more modern than legacy, more flexible than enterprise, more operationally complete than lightweight low-cost tools.",
    ]
    add_two_column_bullets(slide, "Market", left, "Ride Fleet position", right)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.6), Inches(12.2), Inches(5.7)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero.line.color.rgb = PRIMARY_SOFT

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(9.55), Inches(0.95), height=Inches(1.0))

    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.05), Inches(2.15), Inches(0.45)
    )
    label.fill.solid()
    label.fill.fore_color.rgb = PRIMARY
    label.line.color.rgb = PRIMARY
    tf = label.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "TRIANGLE DEALERSHIP"
    run.font.name = "Aptos"
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(255, 255, 255)

    title = slide.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(7.8), Inches(1.6))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Modernize Triangle's loaner program with Ride Fleet"
    r.font.name = "Aptos Display"
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = PRIMARY_DARK

    sub = slide.shapes.add_textbox(Inches(0.95), Inches(3.25), Inches(8.1), Inches(1.0))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Faster service lane operations, cleaner customer communication, stronger billing visibility, and a modern mobile workflow."
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.color.rgb = TEXT

    bullets = [
        "Modern dealership loaner workflow",
        "Billing, accounting, and issue resolution in one platform",
        "Commercial model built to compete with enterprise and legacy vendors",
    ]
    add_bullets(slide, bullets, left=0.95, top=4.35, width=7.6, height=1.4, font_size=17)
    add_footer(slide)


def add_content_slide(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_band(slide, title, subtitle)
    add_bullets(slide, bullets)
    add_footer(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)

    add_content_slide(
        prs,
        "Why Triangle should revisit its loaner stack now",
        "The dealership loaner tool now has to connect service lane, customer, support, and accounting workflows",
        [
            "Customer expectations are mobile-first and speed-sensitive.",
            "Service advisors need faster check-out, return, and exception handling.",
            "Management needs cleaner visibility into utilization, delays, and bottlenecks.",
            "Accounting needs PO support, dealer packets, monthly statements, and cleaner closeout.",
            "Support teams need issue tracking and request-more-info workflows, not scattered follow-up.",
        ],
    )

    add_content_slide(
        prs,
        "What Ride Fleet already does today",
        "This is already live in the platform, not just roadmap language",
        [
            "Dealership loaner intake and borrower packet workflow.",
            "Advisor operations, billing control, and return exception handling.",
            "Extension, vehicle swap, and service completion workflows.",
            "Service lane timeline, billing summary, and accounting closeout.",
            "Printable handoff packet, billing summary, dealer invoice packet, and purchase order print.",
            "Monthly packet and CSV export for accounting and management review.",
            "Alerts for overdue returns, SLA risk, and billing blockers.",
        ],
    )

    add_two_column_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "Service lane strengths",
        [
            "Advisor-facing queue views.",
            "Priority boards for intake, return, billing, and SLA alerts.",
            "Ready-for-pickup visibility.",
            "Borrower packet completion tracking.",
            "Return exception workflow and follow-up.",
        ],
        "Customer and support strengths",
        [
            "Digital customer journey with pre-check-in, agreement, and payment continuity.",
            "Issue and dispute center with status-change communication.",
            "Request-more-info flow with public response link and document upload.",
            "Single issue history and communications timeline.",
        ],
    )
    add_bg(prs.slides[-1])
    add_header_band(prs.slides[-1], "Built for real service lane operations", "Not just reservations and contracts")
    add_footer(prs.slides[-1])

    add_content_slide(
        prs,
        "Accounting and audit control",
        "The platform is designed to help the service lane and accounting team stay aligned",
        [
            "Billing mode and billing status tracking.",
            "Billing contact, authorization reference, and advisor notes.",
            "Accounting closeout with PO number and dealer invoice number support.",
            "Monthly packet and export flow for cleaner internal review.",
            "Printable documents aligned to a customer-ready contract style.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_table_like_slide(slide)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_band(slide, "Where Ride Fleet can win", "Strong modern workflow story with honest boundaries")
    add_two_column_bullets(
        slide,
        "Strong today",
        [
            "Service lane workflow clarity.",
            "Billing and accounting operations.",
            "Issue center and public communications.",
            "Mobile-first user experience.",
            "Fast customization path for dealer-specific workflows.",
        ],
        "Do not overclaim yet",
        [
            "Deep DMS integrations like CDK, Reynolds, or Dealertrack.",
            "OEM-specific forms library across every brand.",
            "Telematics automation for toll, fuel, and mileage capture.",
            "Carrier or insurer network integrations marketed by some enterprise vendors.",
        ],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_pricing_slide(slide)
    add_footer(slide)

    add_content_slide(
        prs,
        "Why this pricing can win",
        "Compete below premium enterprise while staying above commodity expectations",
        [
            "ARSLoaner publicly shows aggressive low-cost pricing, while Dealerware and TSD position higher-value enterprise motions.",
            "Bluebird public pricing is not visible in the materials reviewed.",
            "A simple rooftop plus included-unit model keeps the invoice easy to understand.",
            "A pilot at $199/month and standard at roughly $399/month creates a low-risk first yes.",
            "You preserve room for DMS, telematics, or advanced export work as add-on projects later.",
        ],
    )

    add_content_slide(
        prs,
        "What Triangle should gain",
        "Operational and financial benefits the dealership can feel quickly",
        [
            "Faster check-out and return handling.",
            "Less friction for service advisors and lane coordinators.",
            "Clearer loaner availability and exception management.",
            "Stronger accounting visibility and cleaner monthly close.",
            "Better customer communication and follow-up.",
            "More management visibility into loaner fleet performance and bottlenecks.",
        ],
    )

    add_content_slide(
        prs,
        "Recommended offer for Triangle",
        "A dealer-friendly commercial opening with low switching risk",
        [
            "60-day pilot.",
            "No implementation fee for pilot.",
            "1 rooftop.",
            "Up to 25 active loaners included.",
            "$199/month during pilot.",
            "Success review at day 45 and conversion to standard or group pricing after pilot.",
        ],
    )

    add_content_slide(
        prs,
        "Recommended next step",
        "Use the meeting to get alignment on fit, pilot shape, and must-have integrations",
        [
            "Run a focused workflow demo for intake, packet completion, billing control, extension, swap, and closeout.",
            "Show monthly packet, dealer invoice packet, and issue/request-more-info flow.",
            "Validate their current stack exactly if they reference Bluebird and Rentall together.",
            "Confirm whether DMS integration, OEM forms, or telematics are must-haves for the pilot.",
            "Close on pilot pricing, timeline, and success criteria.",
        ],
    )

    add_content_slide(
        prs,
        "Sources used for this meeting deck",
        "Public materials reviewed on March 24, 2026",
        [
            "Bluebird LoanerTrack: barsnet.com/products/loanertrack and barsnet.com/loanertrack-features",
            "RENTALL dealership loaner: rentallsoftware.com/dealership-loaner",
            "Dealerware loaner fleet and comparison pages: dealerware.com/loaner-fleet and dealerware.com/better-loaner-software",
            "ARSLoaner features and pricing: arsloaner.com/features.aspx and arsloaner.com/pricing.aspx",
            "Reynolds on TSD scale: reyrey.com/company/media-center/news-releases/reynolds-acquires-tsd-mobility-solutions-enhance-fleet",
            "TSD and UVeye announcement: prnewswire.com/news-releases/tsd-mobility-and-uveye-collaborate-to-bring-clarity-and-confidence-to-fleet-management-302710879.html",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build()
